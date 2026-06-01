from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation(r'C:\Users\30810\Desktop\月度汇报\spikingformer_innovation_report.pptx')
title_layout = prs.slide_layouts[0]
FOOTER_COLOR = RGBColor(0x8A, 0x81, 0x74)
TITLE_COLOR = RGBColor(0x17, 0x20, 0x2A)
BODY_COLOR = RGBColor(0x17, 0x20, 0x2A)
SUBTLE_COLOR = RGBColor(0x5E, 0x6F, 0x7E)


def add_slide(prs, title_text, body_lines, footer_text, page_num):
    slide = prs.slides.add_slide(title_layout)
    for ph in slide.placeholders:
        sp = ph._element
        sp.getparent().remove(sp)

    # Title
    txBox = slide.shapes.add_textbox(Cm(1.5), Cm(1.2), Cm(22), Cm(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = 'Microsoft YaHei'
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR

    # Body
    txBox = slide.shapes.add_textbox(Cm(1.5), Cm(4.5), Cm(22), Cm(11))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(body_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = 'Microsoft YaHei'
        p.font.size = Pt(13)
        p.font.color.rgb = BODY_COLOR
        p.space_after = Pt(4)

    # Footer
    txBox = slide.shapes.add_textbox(Cm(1.5), Cm(16.5), Cm(22), Cm(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = footer_text
    p.font.name = 'Microsoft YaHei'
    p.font.size = Pt(8)
    p.font.color.rgb = FOOTER_COLOR

    txBox2 = slide.shapes.add_textbox(Cm(23), Cm(16.5), Cm(2), Cm(1))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = str(page_num).zfill(2)
    p2.font.name = 'Microsoft YaHei'
    p2.font.size = Pt(9)
    p2.font.bold = True
    p2.font.color.rgb = FOOTER_COLOR
    p2.alignment = PP_ALIGN.RIGHT
    return slide


def add_card(slide, x, y, w, h, title, value, unit, desc, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(x + Cm(0.3), y + Cm(0.3), w - Cm(0.6), Cm(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = 'Microsoft YaHei'
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    txBox = slide.shapes.add_textbox(x + Cm(0.3), y + Cm(1.5), w - Cm(0.6), Cm(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = value
    p.font.name = 'Microsoft YaHei'
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    txBox = slide.shapes.add_textbox(x + Cm(0.3), y + Cm(2.5), w - Cm(0.6), Cm(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = unit
    p.font.name = 'Microsoft YaHei'
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    txBox = slide.shapes.add_textbox(x + Cm(0.3), y + Cm(3.0), w - Cm(0.6), Cm(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.name = 'Microsoft YaHei'
    p.font.size = Pt(7)
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


# ========================
# Slide 7: Project overview
# ========================
add_slide(prs,
    'VideoMamba ANN-to-SNN 转换',
    [
        '项目目标：',
        '  将 VideoMamba Small 预训练 ANN 模型转换为 SNN 模型，',
        '  在保持动作识别精度的同时实现脉冲驱动的低功耗推理。',
        '',
        '任务设定：',
        '  数据集：多视角动作视频（12 类），16 帧 @ 224x224',
        '  模型：VideoMamba Small（SSM 架构），~26M 参数',
        '  ANN Baseline：Clean VideoMamba，验证集 Top-1 ~94%',
        '',
        '核心挑战：',
        '  SSM 连续状态 vs 脉冲离散输出之间的表示鸿沟',
        '  ANN 权重向 SNN 的迁移策略（校准 vs 微调 vs 蒸馏）',
        '  LIF 脉冲层阈值设定与深层信息衰减',
    ],
    'VideoMamba / ANN-to-SNN / 月度汇报',
    7
)

# ========================
# Slide 8: Method
# ========================
add_slide(prs,
    'ANN-to-SNN 转换方案',
    [
        '技术路线：',
        '  1. 层结构分析：dump 全部层结构，定位 Mamba Block 内激活层',
        '  2. ReLU -> LIF：将 ReLU 替换为 LIF 脉冲神经元（tau=2.0, 无符号）',
        '  3. 渐进扩展：block0 -> block01 -> block012 -> block0123 -> ...',
        '',
        '脉冲层配置：',
        '  Spike Layer: LIF (Leaky Integrate-and-Fire), backend=torch',
        '  Spike Position: post（在 block 输出后插入 LIF）',
        '  Timesteps: T=4, 无符号 spike {0, 1}',
        '  阈值校准: percentile=0.99, 自动逐层校准',
        '',
        '训练策略：',
        '  知识蒸馏: teacher=ANN best.pth, distill_weight=0.7, T=2.0',
        '  LR: 2e-5 (spike 层 x5), cosine schedule, warmup=1 epoch',
        '  单卡: batch_size=1, update_freq=2, bf16',
    ],
    'VideoMamba / ANN-to-SNN / 月度汇报',
    8
)

# ========================
# Slide 9: Key Results (with cards)
# ========================
slide9 = prs.slides.add_slide(title_layout)
for ph in slide9.placeholders:
    sp = ph._element
    sp.getparent().remove(sp)

# Title
txBox = slide9.shapes.add_textbox(Cm(1.5), Cm(0.6), Cm(22), Cm(2.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = 'LIF SNN 训练结果（Blocks 0,1,2, T=4, 知识蒸馏）'
p.font.name = 'Microsoft YaHei'
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = TITLE_COLOR

# 4 metric cards
card_data = [
    ('ANN\nBaseline', '~94%', 'ANN Top-1', 'Clean VideoMamba\n原始模型', RGBColor(0x14, 0x21, 0x3D)),
    ('LIF SNN\nBlocks 0-2', '91.18%', 'SNN Val Acc', 'Best @ Epoch 2&4\n从 ANN 权重初始化', RGBColor(0x0F, 0x76, 0x6E)),
    ('Initial\n(no train)', '52.21%', 'Initial Acc', '加载权重\n首次前向', RGBColor(0xC7, 0x7A, 0x13)),
    ('Acc Drop\nANN->SNN', '~2.8%', '精度损失', '3 block LIF 转换\n损失可控', RGBColor(0x25, 0x63, 0xEB)),
]

card_w = Cm(5.2)
card_h = Cm(4.2)
start_x = Cm(1.3)
y = Cm(2.8)
gap = Cm(0.45)

for i, (title, value, unit, desc, color) in enumerate(card_data):
    x = start_x + i * (card_w + gap)
    add_card(slide9, x, y, card_w, card_h, title, value, unit, desc, color)

# Training curve (text-based)
txBox = slide9.shapes.add_textbox(Cm(1.5), Cm(7.8), Cm(22), Cm(8))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '训练曲线'
p.font.name = 'Microsoft YaHei'
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = TITLE_COLOR

training_log = [
    'Epoch -1 (init):   Val Loss=1.389,  Val Acc1=52.21%,  Val Acc5=90.44%',
    'Epoch 0:           Train Loss=0.494, Train Acc1=97.47%, Val Acc1=88.24%, Val Acc5=100%',
    'Epoch 1:           Train Loss=0.244, Train Acc1=98.73%, Val Acc1=91.18%, Val Acc5=100%',
    'Epoch 2:           Val Acc1=91.18% (best), Val Acc5=100%',
    'Epoch 3:           Train Loss=0.133, Train Acc1=99.58%, Val Acc1=88.97%',
    'Epoch 4:           Train Loss=0.178, Train Acc1=99.58%, Val Acc1=91.18%',
    '',
    'Spike 验证：三组 LIF 层 spike 输出均为 {0, 1}（无符号，行为正确）',
    '从 52.21% -> 91.18%，知识蒸馏驱动 SNN 成功恢复到接近 ANN 水平',
]
for line in training_log:
    p = tf.add_paragraph()
    p.text = line
    p.font.name = 'Microsoft YaHei'
    p.font.size = Pt(11)
    p.font.color.rgb = SUBTLE_COLOR
    p.space_after = Pt(2)

# Footer
txBox = slide9.shapes.add_textbox(Cm(1.5), Cm(16.5), Cm(22), Cm(1))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = 'VideoMamba / ANN-to-SNN / 月度汇报'
p.font.name = 'Microsoft YaHei'
p.font.size = Pt(8)
p.font.color.rgb = FOOTER_COLOR
txBox2 = slide9.shapes.add_textbox(Cm(23), Cm(16.5), Cm(2), Cm(1))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = '09'
p2.font.name = 'Microsoft YaHei'
p2.font.size = Pt(9)
p2.font.bold = True
p2.font.color.rgb = FOOTER_COLOR
p2.alignment = PP_ALIGN.RIGHT

# ========================
# Slide 10: Ablation & next steps
# ========================
add_slide(prs,
    '消融分析与下一步',
    [
        '已完成消融：No-Train LIF Spike 扫描（block 0-24）',
        '  对各 block 分别插入 LIF 层后直接评估（无训练），发现：',
        '    Block 0-3：性能衰减温和，脉冲层信息保留好',
        '    Block 4+：性能急剧下降，更多 LIF 层 -> 信息丢失加重',
        '  结论：渐进式训练扩展是必要的策略',
        '',
        '当前消融实验矩阵：',
        '  Block 范围：0 -> 0,1 -> 0,1,2 (=91.18%) -> 0-5 (进行中)',
        '  Spike 层位置：post-block（已确认）, pre-block（待验证）',
        '  阈值校准：percentile=0.99, scale 消融待做',
        '',
        '下一步计划：',
        '  1. Blocks 0-5 训练完成并评估（当前 LR=1e-5, 从 0-2 best.pth 继续）',
        '  2. 目标 >=88% 后扩展到 0-7 或 0-11',
        '  3. 超参消融：timesteps (T=2/4/8), threshold_scale, calibration_samples',
        '  4. 多卡 DDP 训练加速',
        '  5. 与 baseline ANN 最终对比（精度-能耗 trade-off）',
    ],
    'VideoMamba / ANN-to-SNN / 月度汇报',
    10
)

# ========================
# Slide 11: Summary
# ========================
slide11 = prs.slides.add_slide(title_layout)
for ph in slide11.placeholders:
    sp = ph._element
    sp.getparent().remove(sp)

txBox = slide11.shapes.add_textbox(Cm(1.5), Cm(1.2), Cm(22), Cm(2.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = 'VideoMamba 项目总结与展望'
p.font.name = 'Microsoft YaHei'
p.font.size = Pt(30)
p.font.bold = True
p.font.color.rgb = TITLE_COLOR

# Left column
txBox = slide11.shapes.add_textbox(Cm(1.5), Cm(4.5), Cm(10.5), Cm(12))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '已完成'
p.font.name = 'Microsoft YaHei'
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = RGBColor(0x0F, 0x76, 0x6E)

done = [
    'ANN Baseline 训练（~94%）',
    '模型结构分析与层定位',
    'ANN-to-SNN 转换脚本完成',
    'LIF blocks 0-2 训练（91.18%）',
    '消融实验框架搭建',
    'Spike 输出验证（{0, 1}）',
    '知识蒸馏训练流程验证',
    'GitHub 同步与 .gitignore 配置',
]
for item in done:
    p = tf.add_paragraph()
    p.text = '  ' + item
    p.font.name = 'Microsoft YaHei'
    p.font.size = Pt(12)
    p.font.color.rgb = BODY_COLOR
    p.space_after = Pt(6)

# Right column
txBox = slide11.shapes.add_textbox(Cm(13.5), Cm(4.5), Cm(10.5), Cm(12))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '待推进'
p.font.name = 'Microsoft YaHei'
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = RGBColor(0xC7, 0x7A, 0x13)

todo = [
    'Blocks 0-5 LIF SNN 训练',
    '更广 block 范围消融',
    '超参数系统化消融',
    '多卡 DDP 训练适配',
    '其他转换方法对比',
    '推理速度与能耗分析',
    '论文图表与报告整理',
]
for item in todo:
    p = tf.add_paragraph()
    p.text = '  ' + item
    p.font.name = 'Microsoft YaHei'
    p.font.size = Pt(12)
    p.font.color.rgb = BODY_COLOR
    p.space_after = Pt(6)

# Footer
txBox = slide11.shapes.add_textbox(Cm(1.5), Cm(16.5), Cm(22), Cm(1))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = 'VideoMamba / ANN-to-SNN / 月度汇报'
p.font.name = 'Microsoft YaHei'
p.font.size = Pt(8)
p.font.color.rgb = FOOTER_COLOR
txBox2 = slide11.shapes.add_textbox(Cm(23), Cm(16.5), Cm(2), Cm(1))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = '11'
p2.font.name = 'Microsoft YaHei'
p2.font.size = Pt(9)
p2.font.bold = True
p2.font.color.rgb = FOOTER_COLOR
p2.alignment = PP_ALIGN.RIGHT

# Save
output_path = r'C:\Users\30810\Desktop\月度汇报\spikingformer_innovation_report_v2.pptx'
prs.save(output_path)
print(f'Done! Total slides: {len(prs.slides)}')
print(f'Saved to: {output_path}')
