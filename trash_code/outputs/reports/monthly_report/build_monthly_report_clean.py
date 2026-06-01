from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"D:\code\PYTHON\video_sm\outputs\monthly_report")
SOURCE_IMAGES = ROOT / "source_slide_images"
OUTPUT = ROOT / "monthly_report_2026_05_clean.pptx"


def cm(value):
    return Inches(value / 2.54)


class C:
    ink = RGBColor(23, 31, 45)
    muted = RGBColor(100, 116, 139)
    line = RGBColor(218, 226, 237)
    paper = RGBColor(255, 255, 255)
    soft = RGBColor(246, 248, 251)
    dark = RGBColor(15, 23, 42)
    dark2 = RGBColor(30, 41, 59)
    blue = RGBColor(37, 99, 235)
    cyan = RGBColor(8, 145, 178)
    green = RGBColor(22, 163, 74)
    amber = RGBColor(217, 119, 6)
    red = RGBColor(220, 38, 38)
    violet = RGBColor(124, 58, 237)


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def line(shape, color=C.line, width=1):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def add_rect(slide, x, y, w, h, color, outline=None, radius=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    fill(shape, color)
    if outline:
        line(shape, outline)
    return shape


def add_text(slide, text, x, y, w, h, size=16, color=C.ink, bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bg(slide, prs, color=C.paper):
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, color)


def add_footer(slide, prs, text="Monthly Research 2026.05"):
    add_rect(slide, cm(1.0), prs.slide_height - cm(0.63), prs.slide_width - cm(2.0), Pt(1), C.line)
    add_text(slide, text, cm(1.0), prs.slide_height - cm(0.48), cm(8), cm(0.25), 7.5, RGBColor(148, 163, 184))


def new_slide(prs, dark=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, C.dark if dark else C.paper)
    return slide


def pill(slide, text, x, y, w, h, color, text_color=C.paper, size=10.5):
    box = add_rect(slide, x, y, w, h, color, radius=True)
    tf = box.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = text_color
    return box


def card(slide, x, y, w, h, title, body, accent=C.blue):
    add_rect(slide, x, y, w, h, C.paper, C.line, radius=True)
    add_rect(slide, x, y, Pt(4), h, accent)
    add_text(slide, title, x + cm(0.35), y + cm(0.28), w - cm(0.7), cm(0.42), 14, C.ink, True)
    body_h = h - cm(1.0)
    if body_h < cm(0.28):
        body_h = cm(0.28)
    add_text(slide, body, x + cm(0.35), y + cm(0.85), w - cm(0.7), body_h, 10.5, C.muted)


def section(slide, label):
    pill(slide, label, cm(1.0), cm(0.55), cm(3.05), cm(0.52), C.dark2, size=10)


def metric(slide, value, label, x, y, color):
    add_rect(slide, x, y, cm(4.15), cm(1.15), C.dark2, RGBColor(51, 65, 85), radius=True)
    add_text(slide, value, x + cm(0.25), y + cm(0.16), cm(2.4), cm(0.42), 20, color, True)
    add_text(slide, label, x + cm(0.25), y + cm(0.72), cm(3.55), cm(0.28), 8.8, RGBColor(203, 213, 225))


def bar_chart(slide, x, y, w, h, labels, values, colors, max_value=100):
    add_rect(slide, x, y + h, w, Pt(1), C.line)
    step = w / len(values)
    for i, value in enumerate(values):
        bw = step * 0.52
        bx = x + step * i + step * 0.24
        bh = h * value / max_value
        by = y + h - bh
        add_rect(slide, bx, by, bw, bh, colors[i])
        add_text(slide, f"{value:.1f}", bx - cm(0.2), by - cm(0.38), bw + cm(0.4), cm(0.25), 8.8, C.ink, True, PP_ALIGN.CENTER)
        add_text(slide, labels[i], bx - cm(0.3), y + h + cm(0.12), bw + cm(0.6), cm(0.55), 7.8, C.muted, False, PP_ALIGN.CENTER)


def mini_table(slide, x, y, w, h, rows, col_ratios, header=C.dark2):
    row_h = h / len(rows)
    total = sum(col_ratios)
    col_ws = [w * r / total for r in col_ratios]
    for r, row in enumerate(rows):
        cx = x
        for c, text in enumerate(row):
            color = header if r == 0 else (C.soft if r % 2 == 0 else C.paper)
            add_rect(slide, cx, y + r * row_h, col_ws[c], row_h, color, C.line)
            txt_color = C.paper if r == 0 else C.ink
            align = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            add_text(slide, str(text), cx + cm(0.13), y + r * row_h + cm(0.12), col_ws[c] - cm(0.26), row_h - cm(0.15),
                     8.4 if r == 0 else 8.1, txt_color, r == 0, align)
            cx += col_ws[c]


def cover(prs):
    slide = new_slide(prs, True)
    add_text(slide, "月度研究汇报", cm(1.0), cm(0.85), cm(8), cm(0.5), 19, RGBColor(203, 213, 225), True)
    add_text(slide, "SpikingFormer 创新点\n与 VideoMamba 脉冲化项目进展", cm(1.0), cm(1.75), cm(14.2), cm(1.7), 31, C.paper, True)
    add_text(slide, "2026.05｜阶段性实验总结与下一步计划", cm(1.03), cm(4.05), cm(9), cm(0.4), 14, RGBColor(203, 213, 225))
    pill(slide, "第一部分：已有创新点 PPT", cm(1.0), cm(5.25), cm(4.9), cm(0.56), C.blue)
    pill(slide, "第二部分：VideoMamba SNN 项目", cm(6.15), cm(5.25), cm(5.4), cm(0.56), C.cyan)
    metric(slide, "94.12%", "ANN val baseline", cm(20.3), cm(1.1), C.green)
    metric(slide, "83.09%", "24-block LIF SNN best", cm(20.3), cm(2.82), C.cyan)
    metric(slide, "24", "active LIF spike layers", cm(20.3), cm(4.54), C.violet)
    add_footer(slide, prs, "Monthly Research 2026.05｜SNN Video Understanding")


def agenda(prs):
    slide = new_slide(prs)
    section(slide, "汇报结构")
    add_text(slide, "本月工作整理为两个相互衔接的创新点", cm(1.0), cm(1.25), cm(20), cm(0.55), 24, C.ink, True)
    add_text(slide, "前半部分保留已有 PPT 的内容；后半部分补充当前 VideoMamba 脉冲化项目的路线、失败尝试、实验结果和下一步计划。", cm(1.05), cm(2.0), cm(21), cm(0.45), 12, C.muted)
    card(slide, cm(1.0), cm(3.0), cm(10.8), cm(3.0), "Part I｜SpikingFormer 创新点",
         "基于粗粒度理解的脉冲信息筛选机制探索：创新动机、整体流程、局部筛选机制，以及失败尝试如何帮助定位问题。", C.blue)
    card(slide, cm(13.0), cm(3.0), cm(10.8), cm(3.0), "Part II｜VideoMamba SNN 项目",
         "以 clean VideoMamba ANN 为起点，比较 ANN2SNN 转换、自定义 signed spike、SpikingJelly LIF 三条路线，形成可训练 unsigned LIF 主线。", C.cyan)
    add_footer(slide, prs)


def part_one(prs):
    slide = new_slide(prs, True)
    add_text(slide, "Part I", cm(1.1), cm(1.0), cm(4), cm(0.45), 18, RGBColor(148, 163, 184), True)
    add_text(slide, "基于粗粒度理解的\n脉冲信息筛选机制探索", cm(1.1), cm(2.0), cm(15), cm(1.4), 31, C.paper, True)
    add_text(slide, "以下 5 页为原 PPT 内容的高清嵌入版，用于保证版式稳定。", cm(1.15), cm(4.0), cm(14), cm(0.4), 13, RGBColor(203, 213, 225))
    pill(slide, "创新动机｜整体流程｜局部筛选机制｜失败尝试的作用", cm(1.15), cm(5.55), cm(11.2), cm(0.56), C.blue)
    add_footer(slide, prs, "Part I｜SpikingFormer Innovation")


def source_slide(prs, image_path):
    slide = new_slide(prs)
    # Keep a small frame so the image never bleeds outside the 16:9 canvas.
    slide.shapes.add_picture(str(image_path), 0, 0, width=prs.slide_width, height=prs.slide_height)


def part_two(prs):
    slide = new_slide(prs, True)
    add_text(slide, "Part II", cm(1.1), cm(1.0), cm(4), cm(0.45), 18, RGBColor(148, 163, 184), True)
    add_text(slide, "VideoMamba 脉冲化项目进展", cm(1.1), cm(2.05), cm(16), cm(0.85), 32, C.paper, True)
    add_text(slide, "从直接转换失败，到基于 ANN 权重的可训练 SNN：当前已完成 full 24-block unsigned LIF 的可行性验证。", cm(1.15), cm(3.35), cm(17.5), cm(0.45), 13, RGBColor(203, 213, 225))
    pill(slide, "ANN baseline 94.12%", cm(1.15), cm(5.2), cm(5.0), cm(0.56), C.green)
    pill(slide, "ANN2SNN direct conversion failed", cm(6.55), cm(5.2), cm(6.1), cm(0.56), C.red)
    pill(slide, "24-block unsigned LIF best 83.09%", cm(13.05), cm(5.2), cm(6.25), cm(0.56), C.cyan)
    add_footer(slide, prs, "Part II｜VideoMamba SNN")


def project_frame(prs):
    slide = new_slide(prs)
    section(slide, "项目定位")
    add_text(slide, "问题不是“把层替换掉”，而是让 VideoMamba 的表示逐步适应脉冲传输", cm(1.0), cm(1.2), cm(22), cm(0.55), 23, C.ink, True)
    card(slide, cm(1.0), cm(2.35), cm(7.05), cm(3.5), "约束",
         "数据集较小，不能从随机初始化学到足够强的时空表示；必须复用 clean ANN 预训练权重。", C.amber)
    card(slide, cm(9.0), cm(2.35), cm(7.05), cm(3.5), "目标",
         "在 Mamba block 外部引入脉冲层，逐步提高脉冲化程度，并确认传输数据尽可能为 {0,1}。", C.cyan)
    card(slide, cm(17.0), cm(2.35), cm(7.05), cm(3.5), "策略",
         "用 teacher distillation 与分阶段扩展 block 范围，缓解一次性插入大量 LIF 导致的精度崩溃。", C.green)
    add_footer(slide, prs)


def conversion_fail(prs):
    slide = new_slide(prs)
    section(slide, "失败路线")
    add_text(slide, "直接 ANN2SNN 转换不适配 VideoMamba", cm(1.0), cm(1.15), cm(20), cm(0.55), 24, C.ink, True)
    add_text(slide, "加载 ANN 权重、插入 spike 层并做阈值校准，但不重新训练。随着脉冲 block 增多，精度快速坍塌。", cm(1.05), cm(1.9), cm(21), cm(0.4), 11.5, C.muted)
    rows = [
        ["run", "blocks", "ANN val", "SNN val", "drop", "test"],
        ["block0", "0", "94.12", "74.26", "-19.85", "72.67"],
        ["block01", "0,1", "94.12", "67.65", "-26.47", "66.15"],
        ["block0123", "0..3", "94.12", "51.47", "-42.65", "55.28"],
    ]
    mini_table(slide, cm(1.0), cm(2.65), cm(11.4), cm(2.25), rows, [1.15, 1.15, 1.1, 1.1, 1.1, 1.1], C.red)
    add_text(slide, "validation acc1", cm(14.2), cm(2.25), cm(5), cm(0.35), 12, C.ink, True)
    bar_chart(slide, cm(14.2), cm(2.85), cm(8.2), cm(2.45), ["ANN", "1 blk", "2 blks", "4 blks"], [94.1, 74.3, 67.6, 51.5],
              [C.green, C.amber, C.amber, C.red])
    card(slide, cm(1.0), cm(5.65), cm(22.0), cm(1.05), "结论",
         "失败路线说明 VideoMamba 的 residual、LayerNorm 与 Mamba 动态不适合简单阈值校准式转换，后续必须转向“加载 ANN 参数 + 插入脉冲层 + 再训练恢复”。", C.red)
    add_footer(slide, prs)


def trainable_route(prs):
    slide = new_slide(prs)
    section(slide, "可训练路线")
    add_text(slide, "当前主线：基于 ANN 参数的可训练 SNN", cm(1.0), cm(1.15), cm(20), cm(0.55), 24, C.ink, True)
    steps = [
        ("1", "加载 ANN", "clean VideoMamba\nbest.pth"),
        ("2", "插入 LIF", "post-block spike\nMamba 不动"),
        ("3", "蒸馏训练", "clean ANN teacher\n约束 logits"),
        ("4", "分阶段扩展", "0..2 → 0..5\n→ 0..11 → 0..23"),
        ("5", "验证输出", "24 层 active LIF\n均为 {0,1}"),
    ]
    for i, (num, title, body) in enumerate(steps):
        x = cm(1.0 + i * 4.7)
        y = cm(2.75)
        add_rect(slide, x, y, cm(4.0), cm(2.0), C.paper, C.line, radius=True)
        pill(slide, num, x + cm(0.25), y + cm(0.23), cm(0.66), cm(0.48), C.cyan, size=10)
        add_text(slide, title, x + cm(1.05), y + cm(0.27), cm(2.6), cm(0.3), 11.5, C.ink, True)
        add_text(slide, body, x + cm(0.25), y + cm(0.9), cm(3.45), cm(0.7), 9.5, C.muted, False, PP_ALIGN.CENTER)
        if i < 4:
            add_rect(slide, x + cm(4.13), y + cm(0.98), cm(0.33), Pt(2), C.line)
    card(slide, cm(1.0), cm(5.55), cm(22.0), cm(1.05), "时间步设置",
         "当前 SNN_TIMESTEPS=4：同一视频输入重复运行 4 次完整 forward，LIF 内部保留膜电位状态，最终平均 4 次 logits。", C.violet)
    add_footer(slide, prs)


def layer_compare(prs):
    slide = new_slide(prs)
    section(slide, "脉冲层对比")
    add_text(slide, "两种脉冲层尝试：精度友好 vs 部署友好", cm(1.0), cm(1.15), cm(20), cm(0.55), 24, C.ink, True)
    rows = [
        ["方案", "输出", "关键特征", "24-block 结果", "判断"],
        ["TrainableSpike3dSeq", "{-θ,0,+θ}", "per-channel threshold", "no-train 77.94；训练 best 96.32", "精度高，但 signed 输出不利于重参数化"],
        ["SpikingJelly LIF", "{0,1}", "MultiStepLIFNode tau=2.0", "initial 25.74；训练 best 83.09", "更符合真正脉冲传输和部署方向"],
    ]
    mini_table(slide, cm(1.0), cm(2.25), cm(22.0), cm(2.25), rows, [1.6, 1.2, 2.2, 2.6, 3.2], C.dark2)
    card(slide, cm(1.0), cm(5.05), cm(10.7), cm(1.45), "signed 层为什么不作为主线",
         "signed 输出带 per-channel threshold scale，虽然更接近 ANN 激活，但不是纯 {0,1} 脉冲，部署和重参数化意义较弱。", C.amber)
    card(slide, cm(12.4), cm(5.05), cm(10.7), cm(1.45), "unsigned LIF 为什么继续做",
         "更难训练，但 spike_stats 已确认 24 个 active LIF 层全部输出 {0,1}，更接近“数据以脉冲形式传输”的目标。", C.cyan)
    add_footer(slide, prs)


def architecture(prs):
    slide = new_slide(prs)
    section(slide, "模型架构")
    add_text(slide, "VideoMamba Unsigned LIF SNN：Mamba block 保持不动，block 输出脉冲化", cm(1.0), cm(1.15), cm(22), cm(0.55), 22, C.ink, True)
    boxes = [
        ("Input video", "B×3×16×224×224", C.blue, 1.0, 2.45, 3.0),
        ("PatchEmbed", "Conv3D\n384×16×14×14", C.cyan, 4.55, 2.45, 3.0),
        ("Tokens + PE", "3137×384\nCLS + pos", C.violet, 8.1, 2.45, 3.0),
        ("24× blocks", "VideoMamba Block_i\n→ LIF_i {0,1}", C.green, 11.65, 2.25, 4.3),
        ("Final norm", "residual + LN", C.amber, 16.55, 2.45, 3.0),
        ("Head", "mean pool\nLinear 12", C.red, 20.1, 2.45, 3.0),
    ]
    for title, body, color, xcm, ycm, wcm in boxes:
        x, y, w = cm(xcm), cm(ycm), cm(wcm)
        h = cm(1.75 if title == "24× blocks" else 1.48)
        add_rect(slide, x, y, w, h, C.paper, color, radius=True)
        add_text(slide, title, x + cm(0.15), y + cm(0.17), w - cm(0.3), cm(0.28), 10.8, color, True, PP_ALIGN.CENTER)
        add_text(slide, body, x + cm(0.15), y + cm(0.62), w - cm(0.3), h - cm(0.65), 8.5, C.muted, False, PP_ALIGN.CENTER)
    for xcm in [4.15, 7.7, 11.25, 16.1, 19.7]:
        add_rect(slide, cm(xcm), cm(3.17), cm(0.22), Pt(2), C.line)
    card(slide, cm(1.0), cm(5.45), cm(7.0), cm(1.1), "插入位置", "每个 Mamba block 后插入 LIF；patch_embed 不做 spike。", C.blue)
    card(slide, cm(8.75), cm(5.45), cm(7.0), cm(1.1), "脉冲输出", "full 24-block 模型中 24 个 active LIF 层均输出 {0,1}。", C.cyan)
    card(slide, cm(16.5), cm(5.45), cm(7.0), cm(1.1), "训练方式", "T=4，双 view，clean ANN teacher distillation。", C.violet)
    add_footer(slide, prs)


def latest_results(prs):
    slide = new_slide(prs)
    section(slide, "最新结果")
    add_text(slide, "Unsigned LIF 分阶段训练：脉冲化程度提高后仍能恢复到可接受精度", cm(1.0), cm(1.15), cm(22), cm(0.55), 22, C.ink, True)
    stages = [
        ("ANN", "94.12", "baseline", C.green),
        ("0..2 LIF", "91.18", "best", C.cyan),
        ("0..5 LIF", "88.97", "best", C.cyan),
        ("0..11 LIF", "78.68", "best", C.amber),
        ("0..23 LIF", "83.09", "best", C.blue),
    ]
    for i, (name, value, note, color) in enumerate(stages):
        x = cm(1.0 + i * 4.55)
        y = cm(2.35)
        add_rect(slide, x, y, cm(3.9), cm(1.45), C.paper, C.line, radius=True)
        add_text(slide, name, x + cm(0.22), y + cm(0.18), cm(3.45), cm(0.28), 10.5, C.muted, True, PP_ALIGN.CENTER)
        add_text(slide, value, x + cm(0.22), y + cm(0.56), cm(3.45), cm(0.45), 19, color, True, PP_ALIGN.CENTER)
        add_text(slide, note, x + cm(0.22), y + cm(1.08), cm(3.45), cm(0.22), 8, C.muted, False, PP_ALIGN.CENTER)
    rows = [
        ("0..2", "initial 52.21 → best 91.18 → latest 91.18"),
        ("0..5", "initial 59.56 → best 88.97 → latest 88.97"),
        ("0..11", "initial 38.97 → best 78.68 → latest 78.68"),
        ("0..23", "initial 25.74 → best 83.09 → latest 81.62"),
    ]
    for i, (stage, text) in enumerate(rows):
        y = cm(4.3 + i * 0.42)
        pill(slide, stage, cm(1.05), y, cm(1.55), cm(0.28), C.dark2, size=7.5)
        add_text(slide, text, cm(2.9), y + cm(0.02), cm(14.5), cm(0.24), 8.8, C.muted)
    card(slide, cm(1.0), cm(6.0), cm(22.0), cm(1.15), "当前判断",
         "full 24-block 一次性初始精度只有 25.74%，但训练 3 个 epoch 已恢复到 83.09%，说明路线可行；后续重点是延长训练、降低学习率微调，并尝试更细粒度插入策略。", C.green)
    add_footer(slide, prs)


def contribution(prs):
    slide = new_slide(prs)
    section(slide, "阶段贡献")
    add_text(slide, "本阶段工作的价值：把“能不能做”推进成“该怎么做”", cm(1.0), cm(1.15), cm(21), cm(0.55), 23, C.ink, True)
    items = [
        ("1", "清理并固定实验工程", "精简项目结构、同步 GitHub、保留训练日志，形成可复现实验流程。", C.blue),
        ("2", "验证直接转换不可行", "ANN2SNN 路线被定量否定，避免继续在错误方向上消耗时间。", C.red),
        ("3", "比较两类脉冲层", "确认 signed spike 精度高但部署意义弱，unsigned LIF 更符合真正脉冲化目标。", C.amber),
        ("4", "建立可训练 SNN 主线", "基于 ANN 预训练、teacher distillation 和分阶段扩展，full 24-block 已达到 83.09%。", C.green),
    ]
    for i, (num, title, body, color) in enumerate(items):
        x = cm(1.0 + (i % 2) * 11.45)
        y = cm(2.25 + (i // 2) * 2.0)
        add_rect(slide, x, y, cm(10.35), cm(1.45), C.paper, C.line, radius=True)
        pill(slide, num, x + cm(0.25), y + cm(0.35), cm(0.62), cm(0.45), color, size=9.5)
        add_text(slide, title, x + cm(1.15), y + cm(0.28), cm(8.5), cm(0.32), 12.5, C.ink, True)
        add_text(slide, body, x + cm(1.15), y + cm(0.78), cm(8.75), cm(0.45), 9.7, C.muted)
    add_footer(slide, prs)


def next_steps(prs):
    slide = new_slide(prs)
    section(slide, "下一步计划")
    add_text(slide, "后续优先把 full 24-block LIF 从“可行”推到“稳定可用”", cm(1.0), cm(1.15), cm(21), cm(0.55), 23, C.ink, True)
    card(slide, cm(1.0), cm(2.3), cm(7.0), cm(3.4), "短期：继续恢复精度",
         "以 24-block best checkpoint 为起点，降低学习率继续训练；观察 val 是否能稳定到 85+ 或 88+。", C.green)
    card(slide, cm(8.8), cm(2.3), cm(7.0), cm(3.4), "中期：细化插入策略",
         "保持 Mamba block 不动，尝试 post-block 归一化、局部分组和阈值初始化，减少分布偏移。", C.cyan)
    card(slide, cm(16.6), cm(2.3), cm(7.0), cm(3.4), "长期：部署可解释性",
         "持续记录 spike_stats：非零率、二值输出、层间分布；再评估能耗、稀疏性与硬件友好性。", C.violet)
    add_footer(slide, prs)


def closing(prs):
    slide = new_slide(prs, True)
    add_text(slide, "阶段结论", cm(1.1), cm(0.95), cm(4.0), cm(0.45), 18, RGBColor(148, 163, 184), True)
    add_text(slide, "本月的核心进展不是一次性完成 SNN，\n而是筛掉了不适配路线，并建立了可继续训练的脉冲化主线。", cm(1.1), cm(1.95), cm(18.5), cm(1.5), 28, C.paper, True)
    add_text(slide, "VideoMamba ANN baseline：94.12% validation acc1\n直接 ANN2SNN：插入 4 个 block 后降到 51.47%，路线暂不作为主线\nfull 24-block unsigned LIF：best 83.09%，已确认 24 层输出 {0,1}",
             cm(1.15), cm(4.15), cm(17.2), cm(1.25), 13, RGBColor(203, 213, 225))
    pill(slide, "下一阶段目标：提升 full 24-block LIF 精度稳定性，并探索更细粒度脉冲化设计", cm(1.15), cm(6.1), cm(14.5), cm(0.56), C.cyan)
    add_footer(slide, prs, "Monthly Research 2026.05｜End")


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)

    cover(prs)
    agenda(prs)
    part_one(prs)
    images = sorted(SOURCE_IMAGES.glob("source_slide_*.png"))
    if not images:
        raise RuntimeError(f"No source slide images found in {SOURCE_IMAGES}")
    for path in images:
        source_slide(prs, path)
    part_two(prs)
    project_frame(prs)
    conversion_fail(prs)
    trainable_route(prs)
    layer_compare(prs)
    architecture(prs)
    latest_results(prs)
    contribution(prs)
    next_steps(prs)
    closing(prs)

    prs.save(OUTPUT)
    print(f"saved={OUTPUT}")
    print(f"slides={len(prs.slides)}")


if __name__ == "__main__":
    build()
