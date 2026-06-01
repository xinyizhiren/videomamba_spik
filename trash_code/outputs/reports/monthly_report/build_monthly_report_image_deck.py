from pathlib import Path
import shutil

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches


ROOT = Path(r"D:\code\PYTHON\video_sm\outputs\monthly_report")
SOURCE_IMAGES = ROOT / "source_slide_images"
SLIDE_DIR = ROOT / "image_deck_slides"
OUTPUT = ROOT / "monthly_report_2026_05_image_stable.pptx"

W, H = 1920, 1080
M = 90


class C:
    dark = (15, 23, 42)
    dark2 = (30, 41, 59)
    ink = (24, 31, 42)
    muted = (100, 116, 139)
    paper = (255, 255, 255)
    soft = (246, 248, 251)
    line = (220, 226, 235)
    blue = (37, 99, 235)
    cyan = (8, 145, 178)
    green = (22, 163, 74)
    amber = (217, 119, 6)
    red = (220, 38, 38)
    violet = (124, 58, 237)
    slate = (203, 213, 225)


FONT_REG = r"C:\Windows\Fonts\msyh.ttc"
FONT_BOLD = r"C:\Windows\Fonts\msyhbd.ttc"


def font(size, bold=False):
    path = FONT_BOLD if bold and Path(FONT_BOLD).exists() else FONT_REG
    return ImageFont.truetype(path, size)


def text_size(draw, text, fnt):
    box = draw.textbbox((0, 0), text, font=fnt)
    return box[2] - box[0], box[3] - box[1]


def wrap_text(draw, text, fnt, max_w):
    lines = []
    for para in text.split("\n"):
        if not para:
            lines.append("")
            continue
        cur = ""
        tokens = []
        buf = ""
        for ch in para:
            if ch == " ":
                if buf:
                    tokens.append(buf)
                    buf = ""
                tokens.append(" ")
            elif ord(ch) < 128:
                buf += ch
            else:
                if buf:
                    tokens.append(buf)
                    buf = ""
                tokens.append(ch)
        if buf:
            tokens.append(buf)
        for tok in tokens:
            test = cur + tok
            if text_size(draw, test, fnt)[0] <= max_w or not cur:
                cur = test
            else:
                lines.append(cur.rstrip())
                cur = tok.lstrip()
        if cur:
            lines.append(cur.rstrip())
    return lines


def draw_text(draw, xy, text, size, fill, bold=False, max_w=None, line_gap=1.15, anchor="la"):
    fnt = font(size, bold)
    x, y = xy
    if max_w:
        lines = wrap_text(draw, text, fnt, max_w)
        line_h = int(size * line_gap)
        for i, line in enumerate(lines):
            draw.text((x, y + i * line_h), line, font=fnt, fill=fill)
        return y + len(lines) * line_h
    draw.text((x, y), text, font=fnt, fill=fill, anchor=anchor)
    return y + size


def round_rect(draw, xy, fill, outline=None, radius=18, width=2):
    draw.rounded_rectangle(xy, radius=radius, fill=fill, outline=outline, width=width)


def pill(draw, xy, text, fill, size=28, text_fill=C.paper):
    x1, y1, x2, y2 = xy
    round_rect(draw, xy, fill, radius=14)
    fnt = font(size, True)
    tw, th = text_size(draw, text, fnt)
    draw.text(((x1 + x2 - tw) / 2, (y1 + y2 - th) / 2 - 3), text, font=fnt, fill=text_fill)


def card(draw, xy, title, body, accent=C.blue):
    x1, y1, x2, y2 = xy
    round_rect(draw, xy, C.paper, C.line, radius=16, width=2)
    draw.rectangle((x1, y1, x1 + 8, y2), fill=accent)
    draw_text(draw, (x1 + 28, y1 + 24), title, 30, C.ink, True, max_w=x2 - x1 - 56)
    draw_text(draw, (x1 + 28, y1 + 78), body, 22, C.muted, False, max_w=x2 - x1 - 56, line_gap=1.35)


def footer(draw, text="Monthly Research 2026.05"):
    draw.line((M, H - 68, W - M, H - 68), fill=C.line, width=2)
    draw_text(draw, (M, H - 48), text, 18, (145, 155, 171))


def base(dark=False):
    im = Image.new("RGB", (W, H), C.dark if dark else C.paper)
    return im, ImageDraw.Draw(im)


def save(im, slides, name):
    path = SLIDE_DIR / f"{len(slides) + 1:02d}_{name}.png"
    im.save(path, quality=95)
    slides.append(path)


def cover(slides):
    im, d = base(True)
    draw_text(d, (M, 70), "月度研究汇报", 56, C.slate, True)
    draw_text(d, (M, 170), "SpikingFormer 创新点", 78, C.paper, True)
    draw_text(d, (M, 285), "与 VideoMamba 脉冲化项目进展", 78, C.paper, True)
    draw_text(d, (M, 430), "2026.05｜阶段性实验总结与下一步计划", 40, C.slate)
    pill(d, (M, 640, 460, 710), "第一部分：已有创新点 PPT", C.blue, 28)
    pill(d, (505, 640, 980, 710), "第二部分：VideoMamba SNN 项目", C.cyan, 28)
    metrics = [("94.12%", "ANN val baseline", C.green), ("83.09%", "24-block LIF SNN best", C.cyan), ("24", "active LIF spike layers", C.violet)]
    for i, (value, label, color) in enumerate(metrics):
        x1, y1 = 1500, 120 + i * 180
        round_rect(d, (x1, y1, 1830, y1 + 120), C.dark2, (58, 73, 94), radius=20, width=3)
        draw_text(d, (x1 + 28, y1 + 20), value, 54, color, True)
        draw_text(d, (x1 + 28, y1 + 78), label, 28, C.slate)
    footer(d, "Monthly Research 2026.05｜SNN Video Understanding")
    save(im, slides, "cover")


def agenda(slides):
    im, d = base()
    pill(d, (M, 55, M + 250, 105), "汇报结构", C.dark2, 22)
    draw_text(d, (M, 155), "本月工作整理为两个相互衔接的创新点", 52, C.ink, True)
    draw_text(d, (M, 235), "前半部分保留已有 PPT 内容；后半部分补充当前 VideoMamba 脉冲化项目的路线、失败尝试、实验结果和下一步计划。", 28, C.muted, max_w=1600)
    card(d, (M, 400, 900, 735), "Part I｜SpikingFormer 创新点", "基于粗粒度理解的脉冲信息筛选机制探索：创新动机、整体流程、局部筛选机制，以及失败尝试如何帮助定位问题。", C.blue)
    card(d, (1020, 400, W - M, 735), "Part II｜VideoMamba SNN 项目", "以 clean VideoMamba ANN 为起点，比较 ANN2SNN 转换、自定义 signed spike、SpikingJelly LIF 三条路线，形成可训练 unsigned LIF 主线。", C.cyan)
    footer(d)
    save(im, slides, "agenda")


def part_one(slides):
    im, d = base(True)
    draw_text(d, (M, 110), "Part I", 42, (148, 163, 184), True)
    draw_text(d, (M, 245), "基于粗粒度理解的", 70, C.paper, True)
    draw_text(d, (M, 345), "脉冲信息筛选机制探索", 70, C.paper, True)
    draw_text(d, (M, 495), "以下 5 页为原 PPT 内容的高清嵌入版，用于保证版式稳定。", 34, C.slate)
    pill(d, (M, 690, 940, 760), "创新动机｜整体流程｜局部筛选机制｜失败尝试的作用", C.blue, 27)
    footer(d, "Part I｜SpikingFormer Innovation")
    save(im, slides, "part_i")


def add_source(slides):
    for src in sorted(SOURCE_IMAGES.glob("source_slide_*.png")):
        im = Image.open(src).convert("RGB").resize((W, H), Image.LANCZOS)
        save(im, slides, src.stem)


def part_two(slides):
    im, d = base(True)
    draw_text(d, (M, 110), "Part II", 42, (148, 163, 184), True)
    draw_text(d, (M, 260), "VideoMamba 脉冲化项目进展", 76, C.paper, True)
    draw_text(d, (M, 390), "从直接转换失败，到基于 ANN 权重的可训练 SNN：当前已完成 full 24-block unsigned LIF 的可行性验证。", 32, C.slate, max_w=1500)
    pill(d, (M, 650, 520, 710), "ANN baseline 94.12%", C.green, 24)
    pill(d, (560, 650, 1130, 710), "ANN2SNN direct conversion failed", C.red, 23)
    pill(d, (1170, 650, 1780, 710), "24-block unsigned LIF best 83.09%", C.cyan, 23)
    footer(d, "Part II｜VideoMamba SNN")
    save(im, slides, "part_ii")


def project_frame(slides):
    im, d = base()
    pill(d, (M, 55, M + 250, 105), "项目定位", C.dark2, 22)
    draw_text(d, (M, 150), "问题不是“把层替换掉”，而是让 VideoMamba 的表示逐步适应脉冲传输", 48, C.ink, True, max_w=1680)
    card(d, (M, 330, 600, 760), "约束", "数据集较小，不能从随机初始化学到足够强的时空表示；必须复用 clean ANN 预训练权重。", C.amber)
    card(d, (660, 330, 1170, 760), "目标", "在 Mamba block 外部引入脉冲层，逐步提高脉冲化程度，并确认传输数据尽可能为 {0,1}。", C.cyan)
    card(d, (1230, 330, W - M, 760), "策略", "用 teacher distillation 与分阶段扩展 block 范围，缓解一次性插入大量 LIF 导致的精度崩溃。", C.green)
    footer(d)
    save(im, slides, "project_frame")


def conversion_fail(slides):
    im, d = base()
    pill(d, (M, 55, M + 250, 105), "失败路线", C.dark2, 22)
    draw_text(d, (M, 150), "直接 ANN2SNN 转换不适配 VideoMamba", 54, C.ink, True)
    draw_text(d, (M, 230), "加载 ANN 权重、插入 spike 层并做阈值校准，但不重新训练。随着脉冲 block 增多，精度快速坍塌。", 28, C.muted, max_w=1550)
    rows = [
        ("run", "blocks", "ANN val", "SNN val", "drop", "test"),
        ("block0", "0", "94.12", "74.26", "-19.85", "72.67"),
        ("block01", "0,1", "94.12", "67.65", "-26.47", "66.15"),
        ("block0123", "0..3", "94.12", "51.47", "-42.65", "55.28"),
    ]
    x0, y0, cw, rh = M, 360, 170, 58
    widths = [180, 160, 150, 150, 150, 150]
    for r, row in enumerate(rows):
        x = x0
        for c, cell in enumerate(row):
            fill = C.red if r == 0 else (C.soft if r % 2 == 0 else C.paper)
            d.rectangle((x, y0 + r * rh, x + widths[c], y0 + (r + 1) * rh), fill=fill, outline=C.line)
            draw_text(d, (x + 16, y0 + r * rh + 15), cell, 22, C.paper if r == 0 else C.ink, r == 0)
            x += widths[c]
    # Simple bars.
    labels = ["ANN", "1 blk", "2 blks", "4 blks"]
    vals = [94.1, 74.3, 67.6, 51.5]
    colors = [C.green, C.amber, C.amber, C.red]
    draw_text(d, (1130, 340), "validation acc1", 28, C.ink, True)
    for i, (lab, val, col) in enumerate(zip(labels, vals, colors)):
        x = 1135 + i * 150
        h = int(val / 100 * 300)
        d.rectangle((x, 710 - h, x + 80, 710), fill=col)
        draw_text(d, (x - 8, 710 - h - 34), f"{val:.1f}", 22, C.ink, True)
        draw_text(d, (x - 5, 725), lab, 20, C.muted)
    card(d, (M, 800, W - M, 930), "结论", "失败路线说明 VideoMamba 的 residual、LayerNorm 与 Mamba 动态不适合简单阈值校准式转换，后续必须转向“加载 ANN 参数 + 插入脉冲层 + 再训练恢复”。", C.red)
    footer(d)
    save(im, slides, "conversion_fail")


def trainable_route(slides):
    im, d = base()
    pill(d, (M, 55, M + 250, 105), "可训练路线", C.dark2, 22)
    draw_text(d, (M, 150), "当前主线：基于 ANN 参数的可训练 SNN", 54, C.ink, True)
    steps = [
        ("1", "加载 ANN", "clean VideoMamba\nbest.pth"),
        ("2", "插入 LIF", "post-block spike\nMamba 不动"),
        ("3", "蒸馏训练", "clean ANN teacher\n约束 logits"),
        ("4", "分阶段扩展", "0..2 → 0..5\n→ 0..11 → 0..23"),
        ("5", "验证输出", "24 层 active LIF\n均为 {0,1}"),
    ]
    for i, (num, title, body) in enumerate(steps):
        x, y = M + i * 350, 360
        round_rect(d, (x, y, x + 280, y + 220), C.paper, C.line, radius=16, width=2)
        pill(d, (x + 24, y + 24, x + 82, y + 78), num, C.cyan, 24)
        draw_text(d, (x + 98, y + 31), title, 26, C.ink, True)
        draw_text(d, (x + 30, y + 120), body, 22, C.muted, max_w=220)
        if i < 4:
            d.line((x + 290, y + 110, x + 332, y + 110), fill=C.line, width=4)
    card(d, (M, 740, W - M, 875), "时间步设置", "当前 SNN_TIMESTEPS=4：同一视频输入重复运行 4 次完整 forward，LIF 内部保留膜电位状态，最终平均 4 次 logits。", C.violet)
    footer(d)
    save(im, slides, "trainable_route")


def layer_compare(slides):
    im, d = base()
    pill(d, (M, 55, M + 250, 105), "脉冲层对比", C.dark2, 22)
    draw_text(d, (M, 150), "两种脉冲层尝试：精度友好 vs 部署友好", 54, C.ink, True)
    card(d, (M, 300, 900, 610), "TrainableSpike3dSeq", "输出：{-θ, 0, +θ}\n阈值：per-channel threshold\n结果：no-train 77.94；训练 best 96.32\n判断：精度高，但 signed 输出不利于重参数化。", C.amber)
    card(d, (1020, 300, W - M, 610), "SpikingJelly MultiStepLIFNode", "输出：{0, 1}\n神经元：tau=2.0, detach_reset=True\n结果：initial 25.74；训练 best 83.09\n判断：更符合真正脉冲传输和部署方向。", C.cyan)
    card(d, (M, 720, W - M, 860), "当前选择", "继续以 unsigned LIF 作为主线。它更难训练，但 spike_stats 已确认 24 个 active LIF 层全部输出 {0,1}，更接近“数据以脉冲形式传输”的目标。", C.green)
    footer(d)
    save(im, slides, "layer_compare")


def architecture(slides):
    im, d = base()
    pill(d, (M, 55, M + 250, 105), "模型架构", C.dark2, 22)
    draw_text(d, (M, 150), "VideoMamba Unsigned LIF SNN：Mamba block 保持不动，block 输出脉冲化", 48, C.ink, True)
    items = [
        ("Input video", "B×3×16×224×224", C.blue),
        ("PatchEmbed", "Conv3D\n384×16×14×14", C.cyan),
        ("Tokens + PE", "3137×384\nCLS + pos", C.violet),
        ("24× blocks", "VideoMamba Block_i\n→ LIF_i {0,1}", C.green),
        ("Final norm", "residual + LN", C.amber),
        ("Head", "mean pool\nLinear 12", C.red),
    ]
    x = M
    for title, body, color in items:
        bw = 255 if title != "24× blocks" else 330
        round_rect(d, (x, 330, x + bw, 500), C.paper, color, radius=16, width=3)
        draw_text(d, (x + 18, 355), title, 24, color, True, max_w=bw - 36)
        draw_text(d, (x + 18, 405), body, 20, C.muted, max_w=bw - 36)
        x_next = x + bw + 55
        if title != "Head":
            d.line((x + bw + 10, 415, x + bw + 45, 415), fill=C.line, width=4)
        x = x_next
    card(d, (M, 680, 600, 825), "插入位置", "每个 Mamba block 后插入 LIF；patch_embed 不做 spike。", C.blue)
    card(d, (660, 680, 1170, 825), "脉冲输出", "full 24-block 模型中 24 个 active LIF 层均输出 {0,1}。", C.cyan)
    card(d, (1230, 680, W - M, 825), "训练方式", "T=4，双 view，clean ANN teacher distillation。", C.violet)
    footer(d)
    save(im, slides, "architecture")


def latest_results(slides):
    im, d = base()
    pill(d, (M, 55, M + 250, 105), "最新结果", C.dark2, 22)
    draw_text(d, (M, 150), "Unsigned LIF 分阶段训练：脉冲化程度提高后仍能恢复到可接受精度", 48, C.ink, True, max_w=1700)
    stages = [
        ("ANN", "94.12", "baseline", C.green),
        ("0..2 LIF", "91.18", "best", C.cyan),
        ("0..5 LIF", "88.97", "best", C.cyan),
        ("0..11 LIF", "78.68", "best", C.amber),
        ("0..23 LIF", "83.09", "best", C.blue),
    ]
    for i, (name, value, note, color) in enumerate(stages):
        x = M + i * 350
        round_rect(d, (x, 310, x + 290, 465), C.paper, C.line, radius=16, width=2)
        draw_text(d, (x + 32, 332), name, 24, C.muted, True)
        draw_text(d, (x + 32, 372), value, 48, color, True)
        draw_text(d, (x + 32, 430), note, 20, C.muted)
    rows = [
        ("0..2", "initial 52.21 → best 91.18 → latest 91.18"),
        ("0..5", "initial 59.56 → best 88.97 → latest 88.97"),
        ("0..11", "initial 38.97 → best 78.68 → latest 78.68"),
        ("0..23", "initial 25.74 → best 83.09 → latest 81.62"),
    ]
    for i, (stage, txt) in enumerate(rows):
        y = 545 + i * 54
        pill(d, (M, y, M + 130, y + 36), stage, C.dark2, 18)
        draw_text(d, (M + 165, y + 4), txt, 25, C.muted)
    card(d, (M, 795, W - M, 925), "当前判断", "full 24-block 一次性初始精度只有 25.74%，但训练 3 个 epoch 已恢复到 83.09%，说明路线可行；后续重点是延长训练、降低学习率微调，并尝试更细粒度插入策略。", C.green)
    footer(d)
    save(im, slides, "latest_results")


def contribution(slides):
    im, d = base()
    pill(d, (M, 55, M + 250, 105), "阶段贡献", C.dark2, 22)
    draw_text(d, (M, 150), "本阶段工作的价值：把“能不能做”推进成“该怎么做”", 50, C.ink, True)
    items = [
        ("1", "清理并固定实验工程", "精简项目结构、同步 GitHub、保留训练日志，形成可复现实验流程。", C.blue),
        ("2", "验证直接转换不可行", "ANN2SNN 路线被定量否定，避免继续在错误方向上消耗时间。", C.red),
        ("3", "比较两类脉冲层", "确认 signed spike 精度高但部署意义弱，unsigned LIF 更符合真正脉冲化目标。", C.amber),
        ("4", "建立可训练 SNN 主线", "基于 ANN 预训练、teacher distillation 和分阶段扩展，full 24-block 已达到 83.09%。", C.green),
    ]
    for i, (num, title, body, color) in enumerate(items):
        x = M + (i % 2) * 870
        y = 320 + (i // 2) * 245
        round_rect(d, (x, y, x + 780, y + 170), C.paper, C.line, radius=16, width=2)
        pill(d, (x + 28, y + 45, x + 86, y + 98), num, color, 22)
        draw_text(d, (x + 115, y + 35), title, 30, C.ink, True)
        draw_text(d, (x + 115, y + 88), body, 22, C.muted, max_w=610)
    footer(d)
    save(im, slides, "contribution")


def next_steps(slides):
    im, d = base()
    pill(d, (M, 55, M + 250, 105), "下一步计划", C.dark2, 22)
    draw_text(d, (M, 150), "后续优先把 full 24-block LIF 从“可行”推到“稳定可用”", 50, C.ink, True)
    card(d, (M, 330, 600, 760), "短期：继续恢复精度", "以 24-block best checkpoint 为起点，降低学习率继续训练；观察 val 是否能稳定到 85+ 或 88+。", C.green)
    card(d, (660, 330, 1170, 760), "中期：细化插入策略", "保持 Mamba block 不动，尝试 post-block 归一化、局部分组和阈值初始化，减少分布偏移。", C.cyan)
    card(d, (1230, 330, W - M, 760), "长期：部署可解释性", "持续记录 spike_stats：非零率、二值输出、层间分布；再评估能耗、稀疏性与硬件友好性。", C.violet)
    footer(d)
    save(im, slides, "next_steps")


def closing(slides):
    im, d = base(True)
    draw_text(d, (M, 110), "阶段结论", 42, (148, 163, 184), True)
    draw_text(d, (M, 245), "本月的核心进展不是一次性完成 SNN，", 58, C.paper, True)
    draw_text(d, (M, 325), "而是筛掉不适配路线，建立可继续训练的脉冲化主线。", 58, C.paper, True)
    draw_text(d, (M, 500), "VideoMamba ANN baseline：94.12% validation acc1", 32, C.slate)
    draw_text(d, (M, 555), "直接 ANN2SNN：插入 4 个 block 后降到 51.47%，路线暂不作为主线", 32, C.slate)
    draw_text(d, (M, 610), "full 24-block unsigned LIF：best 83.09%，已确认 24 层输出 {0,1}", 32, C.slate)
    pill(d, (M, 760, 1320, 825), "下一阶段目标：提升 full 24-block LIF 精度稳定性，并探索更细粒度脉冲化设计", C.cyan, 25)
    footer(d, "Monthly Research 2026.05｜End")
    save(im, slides, "closing")


def build_images():
    if SLIDE_DIR.exists():
        shutil.rmtree(SLIDE_DIR)
    SLIDE_DIR.mkdir(parents=True)
    slides = []
    cover(slides)
    agenda(slides)
    part_one(slides)
    add_source(slides)
    part_two(slides)
    project_frame(slides)
    conversion_fail(slides)
    trainable_route(slides)
    layer_compare(slides)
    architecture(slides)
    latest_results(slides)
    contribution(slides)
    next_steps(slides)
    closing(slides)
    return slides


def build_ppt(slides):
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for image in slides:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(image), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(OUTPUT)


if __name__ == "__main__":
    slide_paths = build_images()
    build_ppt(slide_paths)
    print(f"saved={OUTPUT}")
    print(f"slides={len(slide_paths)}")
