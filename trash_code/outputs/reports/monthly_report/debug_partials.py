from pathlib import Path
import sys

from pptx import Presentation
from pptx.util import Inches

sys.path.insert(0, r"D:\code\PYTHON\video_sm\outputs\monthly_report")
import build_monthly_report_clean as b


out = Path(r"D:\code\PYTHON\video_sm\outputs\monthly_report\partials")
out.mkdir(exist_ok=True)

prs = Presentation()
prs.slide_width = Inches(13.333333)
prs.slide_height = Inches(7.5)

for name, fn in [("cover", b.cover), ("agenda", b.agenda), ("part_one", b.part_one)]:
    fn(prs)
    prs.save(out / f"{len(prs.slides):02d}_{name}.pptx")

for image in sorted(b.SOURCE_IMAGES.glob("source_slide_*.png")):
    b.source_slide(prs, image)
    prs.save(out / f"{len(prs.slides):02d}_source.pptx")

for name, fn in [
    ("part_two", b.part_two),
    ("project_frame", b.project_frame),
    ("conversion_fail", b.conversion_fail),
    ("trainable_route", b.trainable_route),
    ("layer_compare", b.layer_compare),
    ("architecture", b.architecture),
    ("latest_results", b.latest_results),
    ("contribution", b.contribution),
    ("next_steps", b.next_steps),
    ("closing", b.closing),
]:
    fn(prs)
    prs.save(out / f"{len(prs.slides):02d}_{name}.pptx")

print(out)
